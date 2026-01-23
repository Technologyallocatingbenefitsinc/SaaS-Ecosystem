from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import io
import os

# Define Styling Constants
THEMES = {
    "default": {
        "bg_color": (255, 255, 255),
        "title_color": (0, 0, 0),
        "body_color": (80, 80, 80),
        "accent_color": (59, 130, 246) # Blue-500
    },
    "dark": {
        "bg_color": (15, 23, 42),
        "title_color": (255, 255, 255),
        "body_color": (203, 213, 225),
        "accent_color": (139, 92, 246) # Violet-500
    },
    "corporate": {
        "bg_color": (255, 255, 255),
        "title_color": (30, 58, 138),
        "body_color": (71, 85, 105),
        "accent_color": (37, 99, 235) # Blue-600
    },
    "warm": {
        "bg_color": (255, 251, 235),
        "title_color": (120, 53, 15),
        "body_color": (146, 64, 14),
        "accent_color": (217, 119, 6) # Amber-600
    },
    # Aesthetic Themes
    "sunset": {
        "bg_color": (255, 241, 242), # Rose-50
        "title_color": (159, 18, 57), # Rose-800
        "body_color": (136, 19, 55), # Rose-900
        "accent_color": (251, 113, 133) # Rose-400
    },
    "forest": {
        "bg_color": (240, 253, 244), # Green-50
        "title_color": (20, 83, 45), # Green-800
        "body_color": (22, 101, 52), # Green-700
        "accent_color": (22, 163, 74) # Green-600
    },
    "ocean": {
        "bg_color": (236, 254, 255), # Cyan-50
        "title_color": (12, 74, 110), # Sky-900
        "body_color": (8, 51, 68), # Cyan-950
        "accent_color": (6, 182, 212) # Cyan-500
    },
    "luxury": {
        "bg_color": (24, 24, 27), # Zinc-900
        "title_color": (234, 179, 8), # Yellow-500 (Gold)
        "body_color": (228, 228, 231), # Zinc-200
        "accent_color": (250, 204, 21) # Yellow-400
    }
}

def draw_theme_layout(slide, theme_name, theme_colors, width, height):
    """Adds geometric shapes to creates a modern layout."""
    shapes = slide.shapes
    accent_rgb = RGBColor(*theme_colors["accent_color"])
    
    if theme_name == "corporate":
        # Left Accent Bar
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.5), height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = accent_rgb
        shape.line.fill.background() # No outline
    
    elif theme_name == "dark":
        # Top thin accent line
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, Inches(0.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = accent_rgb
        shape.line.fill.background()
        
        # Bottom right accent circle
        circle = shapes.add_shape(MSO_SHAPE.OVAL, width - Inches(2), height - Inches(2), Inches(3), Inches(3))
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent_rgb
        circle.fill.transparency = 0.8 # Fades out
        circle.line.fill.background()

    elif theme_name == "warm":
        # Soft top header block
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, Inches(1.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(253, 230, 138) # Light yellow accent
        shape.line.fill.background()

    elif theme_name == "sunset":
        # Bottom accent bar
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, height - Inches(1), width, Inches(1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = accent_rgb
        shape.line.fill.background()
        
    elif theme_name == "forest":
        # Left sidebar thin
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.8), height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = accent_rgb
        shape.line.fill.background()

    elif theme_name == "ocean":
        # Top wave accent
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = accent_rgb
        shape.line.fill.background()
        
    elif theme_name == "luxury":
        # Gold Frame
        # Top
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, Inches(0.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = accent_rgb
        shape.line.fill.background()
        # Bottom
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, height - Inches(0.1), width, Inches(0.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = accent_rgb
        shape.line.fill.background()

def generate_pptx(slide_data: list, watermark: bool = False, theme_name: str = "default", aspect_ratio: str = "16:9") -> bytes:
    """
    Generates a style-aware PPTX file interactively.
    """
    prs = Presentation()
    
    # Set Slide Size
    if aspect_ratio == "1:1": # Square Carousel
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(10)
    else: # Default 16:9
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    # Get Theme Colors
    theme = THEMES.get(theme_name, THEMES["default"])
    bg_rgb = RGBColor(*theme["bg_color"])
    title_rgb = RGBColor(*theme["title_color"])
    body_rgb = RGBColor(*theme["body_color"])
    
    for i, slide_info in enumerate(slide_data):
        title_text = slide_info.get("title", "Untitled")
        content_raw = slide_info.get("content", "")
        points = slide_info.get("points", [])
        
        # Create Slide (Blank layout gives us full control, but Title+Content is easier for placeholders)
        # using Blank (6) to allow custom shape drawing easily? 
        # No, let's stick to 1 but adjust positions if needed.
        # Actually, for "Modern", we often want custom text boxes.
        # Let's stick to layout 1 for robustness but style it heavily.
        slide_layout = prs.slide_layouts[1] 
        slide = prs.slides.add_slide(slide_layout)
        
        # --- Apply Background ---
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_rgb
        
        # --- Draw Theme Geometry (At the back) ---
        draw_theme_layout(slide, theme_name, theme, prs.slide_width, prs.slide_height)

        # --- Set Title ---
        title = slide.shapes.title
        if title:
            title.text = title_text
            # Use 'warm' theme offset if header bar exists
            if theme_name == "warm":
                title.top = Inches(0.2)
            
            # Style Title
            for paragraph in title.text_frame.paragraphs:
                paragraph.font.color.rgb = title_rgb
                paragraph.font.bold = True
                paragraph.font.name = "Arial" # Safer font

        # --- Set Content ---
        if len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            
            # Adjust body position for Corporate layout (shift right past bar)
            if theme_name == "corporate":
                body_shape.left = Inches(1.0)
                body_shape.width = prs.slide_width - Inches(1.5)

            tf = body_shape.text_frame
            tf.clear() # Clear existing placeholder text
            
            if isinstance(points, list) and points:
                # Add text
                for i, point in enumerate(points):
                    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                    p.text = point
                    p.level = 0
                    p.font.color.rgb = body_rgb
                    p.font.size = Pt(24) # Start size

                # Auto-fit Logic: Reduce font size if overflowing
                # Heuristic: Check rough character count vs box size
                # Better: pptx doesn't give us calculated h, so we just safeguard with smaller fonts for long lists
                total_chars = sum(len(p) for p in points)
                if total_chars > 300:
                    for p in tf.paragraphs:
                        p.font.size = Pt(18)
                elif total_chars > 500:
                    for p in tf.paragraphs:
                        p.font.size = Pt(14)
                
            else:
                tf.text = str(content_raw)
                for paragraph in tf.paragraphs:
                    paragraph.font.color.rgb = body_rgb
                    paragraph.font.size = Pt(24)
                
                # Auto-fit for block text
                if len(str(content_raw)) > 400:
                    for p in tf.paragraphs:
                        p.font.size = Pt(18)
                if len(str(content_raw)) > 700:
                    for p in tf.paragraphs:
                        p.font.size = Pt(14)

        # --- Add Image (if present) ---
        image_url = slide_info.get("image_url")
        if image_url:
            try:
                # Resolve path
                image_path = None
                if image_url.startswith("/uploads/"):
                   # Map /uploads/xyz -> user_uploads/xyz
                   image_path = image_url.replace("/uploads/", "user_uploads/", 1)
                
                if image_path and os.path.exists(image_path):
                    # add_picture(image_file, left, top, width, height)
                    # Place it on the right key side or bottom
                    # Let's put it on the right side, shrinking text box?
                    # For simplicity, let's put it in the bottom right corner, overlaying if needed
                    # Or adjust the layout.
                    
                    # Modern Layout: Image on right half
                    # 1. Shrink text box
                    if len(slide.placeholders) > 1:
                        body_shape = slide.placeholders[1]
                        body_shape.width = Inches(6.5) # Half width approx (Slide is 13.33)
                    
                    # 2. Add Image
                    img_left = Inches(7)
                    img_top = Inches(2)
                    img_height = Inches(4.5)
                    # Auto width
                    slide.shapes.add_picture(image_path, img_left, img_top, height=img_height)
            except Exception as e:
                print(f"Failed to add image: {e}")

                
        if watermark:
            # Add Discrete Watermark to EACH slide
            # Bottom left corner
            txBox = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.4), Inches(4), Inches(0.3))
            tf = txBox.text_frame
            tf.text = "Generated by MODYFIRE"
            p = tf.paragraphs[0]
            p.font.size = Pt(10)
            p.font.italic = True
            p.font.color.rgb = RGBColor(128, 128, 128) # Gray

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()
